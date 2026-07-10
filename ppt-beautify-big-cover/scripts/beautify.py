#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ppt-beautify: PPT 原地换肤美化 (in-place re-skin)
核心：打开原 pptx，遍历每个 shape 原地改色/字体；递归进入 GROUP 处理子元素。
绝不重建幻灯片（add_group_shape 复制 GROUP 会错位 -> 文字飞出画面）。零坐标误差。
"""
import sys
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------- 主题 ----------------
THEMES = {
    "blue-gold": {  # 深蓝 + 金（党政 / 思政 / 教育，庄重）
        "bg": RGBColor(0x0E, 0x2A, 0x47),
        "card": RGBColor(0x16, 0x3B, 0x5C),
        "accent": RGBColor(0xC8, 0x90, 0x2E),
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),
        "text_blue": RGBColor(0xBF, 0xD3, 0xE6),
        "footer": RGBColor(0x8A, 0xA0, 0xB5),
        "mode": "dark",
    },
    "red-gold": {  # 中国红 + 金（主旋律）
        "bg": RGBColor(0x7A, 0x14, 0x18),
        "card": RGBColor(0x9A, 0x25, 0x2A),
        "accent": RGBColor(0xD4, 0xAF, 0x37),
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),
        "text_blue": RGBColor(0xF0, 0xD9, 0xD9),
        "footer": RGBColor(0xC9, 0xA9, 0xA9),
        "mode": "dark",
    },
    "teal": {  # 科技青
        "bg": RGBColor(0x0B, 0x2E, 0x3A),
        "card": RGBColor(0x11, 0x4B, 0x5F),
        "accent": RGBColor(0x2E, 0xC4, 0xB6),
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),
        "text_blue": RGBColor(0xC5, 0xE3, 0xE8),
        "footer": RGBColor(0x9A, 0xB8, 0xBE),
        "mode": "dark",
    },
}


def is_group(sh):
    try:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            return True
    except Exception:
        pass
    return "GROUP" in str(sh.shape_type)


def is_picture(sh):
    try:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    except Exception:
        pass
    return "PICTURE" in str(sh.shape_type)


def is_auto_shape(sh):
    try:
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return True
    except Exception:
        pass
    return "AUTO_SHAPE" in str(sh.shape_type)


def get_run_color(run):
    try:
        return run.font.color.rgb
    except Exception:
        return None


def is_loud(rgb):
    """红 / 橙等刺眼色 -> 收为强调色"""
    if rgb is None:
        return False
    r, g, b = rgb[0], rgb[1], rgb[2]
    return (r > 150) and (g < 130) and (b < 130)


def set_run_color(run, rgb):
    """设置 run 文字色；若 python-pptx 因"继承/方案色"设置失败，直接写 XML 兜底。"""
    try:
        run.font.color.rgb = rgb
        return
    except Exception:
        pass
    try:
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:schemeClr",
                    "a:srgbClr", "a:prstClr", "a:patternFill"):
            for el in rPr.findall(".//" + qn(tag)):
                el.getparent().remove(el)
        sf = rPr.makeelement(qn("a:solidFill"), {})
        c = sf.makeelement(qn("a:srgbClr"),
                           {"val": "%02X%02X%02X" % (rgb[0], rgb[1], rgb[2])})
        sf.append(c)
        rPr.append(sf)
    except Exception:
        pass


def choose_text_color(run, theme):
    orig = get_run_color(run)
    if is_loud(orig):
        return theme["accent"]
    size = run.font.size.pt if run.font.size else 18
    bold = run.font.bold
    if bold or size >= 24:
        return theme["text_white"]
    return theme["text_blue"]


def recolor_text(shape, theme):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            set_run_color(run, choose_text_color(run, theme))


def recolor_shape(shape, theme):
    if is_picture(shape):
        try:
            shape.line.color.rgb = theme["accent"]
            shape.line.width = Pt(1.5)
        except Exception:
            pass
    elif is_auto_shape(shape):
        # 全量换肤：纯色形状统一染成主题卡片色 + 强调细边
        # （v4 验证：深红/橙等异色卡片也一并收编，视觉才统一）
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = theme["card"]
        except Exception:
            pass
        try:
            shape.line.color.rgb = theme["accent"]
            shape.line.width = Pt(1.25)
        except Exception:
            pass


def walk(shapes, theme, fn):
    for sh in shapes:
        fn(sh, theme)
        if is_group(sh):
            try:
                walk(sh.shapes, theme, fn)
            except Exception:
                pass


# ---------------- 封面大图模式辅助 ----------------
def _sp_tree(slide):
    return slide.shapes._spTree


def send_to_back(slide, shape):
    tree = _sp_tree(slide)
    el = shape._element
    tree.remove(el)
    tree.insert(2, el)  # grpSpPr 之后，置于最底层


def bring_to_front(slide, shape):
    tree = _sp_tree(slide)
    el = shape._element
    tree.remove(el)
    tree.append(el)


def set_solid_alpha(shape, rgb, alpha):
    """纯色填充 + 透明度 alpha(0-100000, 60000=60%)"""
    spPr = shape.fill._xPr
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill",
                "a:grpFill", "a:pattFill"):
        for el in spPr.findall(".//" + qn(tag)):
            el.getparent().remove(el)
    sf = spPr.makeelement(qn("a:solidFill"), {})
    c = sf.makeelement(qn("a:srgbClr"),
                       {"val": "%02X%02X%02X" % (rgb[0], rgb[1], rgb[2])})
    c.append(c.makeelement(qn("a:alpha"), {"val": str(int(alpha))}))
    sf.append(c)
    spPr.append(sf)


def cover_hero(slide, theme, sw, sh):
    """封面大图：选最大图片铺满全屏，删掉封面所有图片/分组/全屏块，
    用该图原始数据重铺一张全屏 hero，加半透明深色蒙版，标题置顶。"""
    from io import BytesIO
    # 递归收集所有图片，选面积最大者
    pics = []

    def collect(shapes):
        for sh in shapes:
            if is_picture(sh):
                pics.append(sh)
            elif is_group(sh):
                collect(sh.shapes)

    collect(slide.shapes)
    if not pics:
        return
    hero = max(pics, key=lambda p: int(p.width) * int(p.height))
    try:
        blob = hero.image.blob
    except Exception:
        blob = None
    tree = _sp_tree(slide)
    # 删除封面所有图片、分组、全屏形状（仅作用于输出文件，原稿不动）
    for s in list(slide.shapes):
        if is_picture(s):
            try:
                tree.remove(s._element)
            except Exception:
                pass
        elif is_group(s):
            try:
                tree.remove(s._element)
            except Exception:
                pass
        elif is_auto_shape(s) and s.width >= sw * 0.98 and s.height >= sh * 0.98:
            try:
                tree.remove(s._element)
            except Exception:
                pass
    # 用原图数据重铺一张全屏 hero
    if blob:
        try:
            nh = slide.shapes.add_picture(BytesIO(blob), 0, 0, sw, sh)
            send_to_back(slide, nh)
        except Exception:
            pass
    # 半透明深色蒙版，保证标题可读
    mask = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    set_solid_alpha(mask, theme["bg"], 55000)
    try:
        mask.line.fill.background()
    except Exception:
        pass
    # 文字置顶（压在蒙版之上）
    for s in list(slide.shapes):
        if s.has_text_frame and s.text_frame.text.strip():
            bring_to_front(slide, s)


def beautify(input_path, output_path, theme_name, add_footer=False, footer_text="演示文档", cover="keep", cover_slide=1):
    theme = THEMES[theme_name]
    prs = Presentation(input_path)
    SW, SH = prs.slide_width, prs.slide_height
    total = len(prs.slides)

    for idx, slide in enumerate(prs.slides, 1):
        # 封面大图模式：单图铺满 + 蒙版 + 标题置顶
        if cover == "hero" and idx == cover_slide:
            cover_hero(slide, theme, SW, SH)
        # 1) 背景
        try:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = theme["bg"]
        except Exception:
            pass

        # 2) 遍历所有形状（含递归分组）
        def shape_fn(sh, t):
            if sh.width >= SW * 0.98 and sh.height >= SH * 0.98 and is_auto_shape(sh):
                # 全屏形状当背景，不染成卡片色
                try:
                    sh.fill.solid()
                    sh.fill.fore_color.rgb = t["bg"]
                except Exception:
                    pass
                return
            recolor_shape(sh, t)
            recolor_text(sh, t)
            # 表格单元格（python-pptx 不通过 has_text_frame 暴露，需单独处理）
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        recolor_text(cell, t)

        walk(slide.shapes, theme, shape_fn)

        # 3) 可选页脚 + 页码
        if add_footer:
            tb = slide.shapes.add_textbox(
                Inches(0.5), SH - Inches(0.42), Inches(8), Inches(0.3)
            )
            tf = tb.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = footer_text
            r.font.size = Pt(10)
            r.font.color.rgb = theme["footer"]

            pn = slide.shapes.add_textbox(
                SW - Inches(2.2), SH - Inches(0.42), Inches(1.7), Inches(0.3)
            )
            pf = pn.text_frame
            pp = pf.paragraphs[0]
            pp.alignment = 2  # 右对齐
            rr = pp.add_run()
            rr.text = "%02d / %02d" % (idx, total)
            rr.font.size = Pt(10)
            rr.font.color.rgb = theme["footer"]

    prs.save(output_path)
    print("OK -> %s  (theme=%s, slides=%d)" % (output_path, theme_name, total))


def main():
    ap = argparse.ArgumentParser(description="PPT 原地换肤美化")
    ap.add_argument("input", help="输入 pptx 路径")
    ap.add_argument("output", help="输出 pptx 路径")
    ap.add_argument("--theme", default="blue-gold", choices=list(THEMES.keys()))
    ap.add_argument("--footer", action="store_true", help="添加页脚与页码")
    ap.add_argument("--footer-text", default="演示文档", help="页脚文字")
    ap.add_argument("--cover", default="keep", choices=["keep", "hero"],
                    help="封面模式: keep=保留原版式; hero=单张大图铺满(电影海报感)")
    ap.add_argument("--cover-slide", type=int, default=1, help="启用 hero 时的封面页序号")
    args = ap.parse_args()
    beautify(args.input, args.output, args.theme, args.footer,
             args.footer_text, args.cover, args.cover_slide)


if __name__ == "__main__":
    main()
