#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
qa_color_audit.py —— PPT 换肤后的程序化颜色审计（无多模态也能验证）。
递归扫描所有 run 文字色 + 形状填充色，报告每页背景色、文字色分布、卡片填充色，
并标出是否还有红/橙等刺眼色残留、图片是否丢失。
"""
import sys
from collections import Counter
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE


def is_group(sh):
    try:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            return True
    except Exception:
        pass
    return "GROUP" in str(sh.shape_type)


def is_picture(sh):
    try:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    except Exception:
        pass
    return "PICTURE" in str(sh.shape_type)


def get_run_color(run):
    try:
        return run.font.color.rgb
    except Exception:
        return None


def is_loud(rgb):
    """rgb 为 'RRGGBB' 十六进制字符串（RGBColor 的 str 形式）"""
    if rgb is None:
        return False
    try:
        r = int(rgb[0:2], 16)
        g = int(rgb[2:4], 16)
        b = int(rgb[4:6], 16)
    except Exception:
        return False
    return (r > 150) and (g < 130) and (b < 130)


def walk(shapes, runs, fills, pics):
    for sh in shapes:
        if is_picture(sh):
            pics.append(1)
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    c = get_run_color(r)
                    if c is not None:
                        runs.append(str(c))
        try:
            if sh.fill.type is not None and "SOLID" in str(sh.fill.type):
                fills.append(str(sh.fill.fore_color.rgb))
        except Exception:
            pass
        if is_group(sh):
            walk(sh.shapes, runs, fills, pics)


def main():
    if len(sys.argv) < 2:
        print("usage: qa_color_audit.py file.pptx")
        sys.exit(1)
    prs = Presentation(sys.argv[1])
    total_pics = 0
    print("=== 颜色审计: %s ===" % sys.argv[1])
    for i, slide in enumerate(prs.slides, 1):
        runs, fills, pics = [], [], []
        bg = "?"
        try:
            srgb = slide.background._element.findall(".//" + qn("a:srgbClr"))
            bg = srgb[0].get("val") if srgb else "none"
        except Exception:
            bg = "err"
        walk(slide.shapes, runs, fills, pics)
        loud = [c for c in runs if is_loud(c)]
        total_pics += len(pics)
        print(
            "slide%d: bg=%s | text=%s"
            % (i, bg, dict(Counter(runs)))
        )
        print(
            "        fills=%s | pics=%d | loud_residue=%s"
            % (dict(Counter(fills)), len(pics), loud if loud else "NONE")
        )
    print("\n总图片数: %d" % total_pics)
    print("提示: 内容完整性请用 markitdown 对比 原文/改后 文本。")


if __name__ == "__main__":
    main()
